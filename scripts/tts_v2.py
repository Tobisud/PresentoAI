#####set up django enviroment
import os
import sys
project_path = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
sys.path.append(project_path)
# Set up Django settings
os.environ.setdefault("DJANGO_SETTINGS_MODULE", "presentoapi.settings")
# Initialize Django
import django
django.setup()
from presento.views import upload_proc_percentage as cur_per
######

###libraries

from pydub import AudioSegment
from moviepy.editor import ImageClip, concatenate_videoclips
from moviepy.editor import VideoFileClip, AudioFileClip
import shutil
import subprocess
from comtypes import client
from pptx import Presentation
from TTS.utils.manage import ModelManager
from TTS.utils.synthesizer import Synthesizer
import argparse
import unicodedata
import logging
import pyttsx3
from close_capture import transcribe as scribe
from close_capture import generate_subtitle_file as cc

class presentoai():
    def __init__(self, procID, voice_model):
        self.procID = procID
        self.voice_model=voice_model
        self.process_percentage = 0

    @staticmethod
    def sort_key(file_name):
        # Extract the numeric part of the file name for sorting
        base_name, _ = os.path.splitext(file_name)
        return int(base_name)

    @staticmethod
    def preprocess_text(text):
        """
        Remove unsupported characters and special characters from the input text.
        """
        # Normalize the text to remove combining characters
        text = ''.join(c for c in text if not unicodedata.combining(c))
        
        # Replace characters not supported by the target encoding
        # Here we use ASCII for simplicity, replace unsupported characters with an empty string
        text = text.encode('ascii', 'replace').decode('ascii')      
        return text
    
    # Method to update process percentage
    def update_process_percentage(self, percentage): 
        return  cur_per(self.procID,percentage)

    # Method to retrieve the current process percentage
    def get_process_percentage(self):
        return self.process_percentage

    #get speaker not from slides
    def getText(self, file_path):
        ppt = Presentation(file_path)
        notes = []
        for slide in ppt.slides:
            text_note = slide.notes_slide.notes_text_frame.text
            notes.append(self.preprocess_text(text_note))  # Preprocess text here
        return notes


    #convert speaker note into wav files
    def getWav (self,folder_path, model_choice):
        folder_name=os.path.basename(folder_path)
        file_path = os.path.abspath(os.path.join(folder_path,folder_name+'.pptx'))
        des_path = os.path.join(folder_path,"To_Wavs")
        if not os.path.exists(des_path):
            os.makedirs(des_path)
        contents = self.getText(file_path)
        print(model_choice)
        match int(model_choice):
            case 1:
                models_path = os.path.join(os.getcwd(), 'TTS','.models.json')
                print(models_path)
                model_manager = ModelManager(models_path)
                model_path, config_path, model_item = model_manager.download_model("tts_models/en/ljspeech/tacotron2-DDC_ph")
                voc_path, voc_config_path, _ = model_manager.download_model(model_item["default_vocoder"])
                syn = Synthesizer(
                    tts_checkpoint=model_path,
                    tts_config_path=config_path,
                    vocoder_checkpoint=voc_path,
                    vocoder_config=voc_config_path
                )
                #create wav files
                for i,text in enumerate(contents): 
                    # outputs = syn.tts(text)
                    # syn.save_wav(outputs, os.path.join(des_path, f"{i+1}.wav"))  
                    # text = self.preprocess_text(text)  # Preprocess text here
                    try:
                        outputs = syn.tts(text)
                        syn.save_wav(outputs, os.path.join(des_path, f"{i + 1}.wav"))
                    except Exception as e:
                        logging.error(f"Error generating or saving WAV for text: {text} - {str(e)}")
                        print(f"Error processing text {i + 1}: {text}")           
            case 2: #using pytts to create male voice (rate -10)
                engine = pyttsx3.init()
                rate = engine.getProperty('rate')
                engine.setProperty('rate', rate-10)
                for i,text in enumerate(contents): 
                    mytext = text 
                    engine.save_to_file(mytext, os.path.join(des_path, str(i+1)+".wav"))
                    engine.runAndWait()        
        
    #cobine all wavs file into a complete wav
    def combineAudio (self, file_path):
        audio_path=os.path.join(file_path,"To_Wavs")
        des_path=os.path.join(file_path,"To_Wavs")
        files = os.listdir(audio_path)
        # Filter for .mp3 files
        wav_files = [file for file in files if file.endswith('.wav')]
        # Sort the files if needed 
        wav_files=sorted(wav_files,key=self.sort_key)
        # Initialize an empty AudioSegment
        full_audio = AudioSegment.silent(500)
        # Process and concatenate each .mp3 file
        for wav_file in wav_files:
            file_path = os.path.join(audio_path, wav_file)
            audio = AudioSegment.from_wav(file_path)
            # Concatenate the audio file & add delay 1s in between each audio
            full_audio += audio + AudioSegment.silent(500)
            print("Successfully combine an audio" + wav_file)
        # Export the concatenated audio to a new file
        output_path=os.path.join(des_path,"full_audio.wav")
        full_audio.export(output_path, format='wav')
        print("Successfully combine all audio" + des_path+"\\"+"full_audio.wav")

    def getCC (self,file_path):
        audio_path=os.path.join(file_path,"To_Wavs","full_audio.wav")
        language, segments = scribe(audio_path)
        subtitle_file = cc(language,segments,file_path)
        return subtitle_file

    #get duration from each wav files
    def getDuration(self,file_path):
        des_path=os.path.join(file_path,"To_Wavs")
        files = os.listdir(des_path)
        # Filter for .mp3 files
        wav_files = [file for file in files if file.endswith('.wav')]
        # Sort the files if needed
        wav_files=sorted(wav_files,key=presentoai.sort_key)
        duration=[]
        for wav_file in wav_files:
            file_path = os.path.join(des_path, wav_file)
            audio = AudioSegment.from_wav(file_path)
            duration.append((audio.duration_seconds)+0.5)
        return duration

    #convert pptx to png files
    def getImg (self,folder_path):
        folder_name=os.path.basename(folder_path)
        file_path = os.path.abspath(os.path.join(folder_path,folder_name+'.pptx'))
        des_path = os.path.join(folder_path,"To_Pngs")
        if not os.path.exists(des_path):
            os.makedirs(des_path)
        powerpoint = client.CreateObject('Powerpoint.Application')
        presentation=powerpoint.Presentations.Open(file_path)
        slides_count = presentation.Slides.Count
        for i in range(1, slides_count + 1):
            slide = presentation.Slides(i)
            slide.Export(os.path.join(des_path, f"{i}.png"), "PNG")
        powerpoint.ActivePresentation.Close()
        powerpoint.Quit()


    #using png files to create clip
    def getClip(self,file_path):
        folder_path = os.path.join(file_path,"To_Pngs")
        # Define the durations for each image clip (in seconds)
        des_path=os.path.join(file_path, "To_Mp4")
        if not os.path.exists(des_path):
            os.makedirs(des_path)
        durations = self.getDuration(file_path)  # get duration of each audio
        # List all files in the folder
        files = os.listdir(folder_path)
        # Filter for image files (assuming .jpg files, you can adjust this as needed)
        image_files = [file for file in files if file.endswith('.png')]
        # Sort the image files if needed 
        image_files=sorted(image_files,key=self.sort_key)
        # Check if the number of durations matches the number of image files
        if len(durations) != len(image_files):
            raise ValueError("The number of durations must match the number of image files")
        # Create video clips from each image with the specified durations
        clips = []
        for image_file, duration in zip(image_files, durations):
            file_path = os.path.join(folder_path, image_file)
            clip = ImageClip(file_path).set_duration(duration)
            clips.append(clip)
        # Concatenate the clips
        temp_clip = concatenate_videoclips(clips)
        # Define the output video file path
        output_path = os.path.join(des_path,"temp.mp4")
        # Export the final video
        temp_clip.write_videofile(output_path, codec='libx264', fps=1)

    #combine the full audio file with the clip
    def video_wsound(self,file_path):
        #define
        video_path = os.path.join(file_path, "To_Mp4", "temp.mp4")
        audio_path = os.path.join(file_path, "To_Wavs", "full_audio.wav")
        output_folder = os.path.join(file_path, "temp")
        if not os.path.exists(output_folder):
            os.makedirs(output_folder)
        output_path = os.path.join(output_folder,self.procID+".mp4")
        # Load the video file
        video_clip = VideoFileClip(video_path)
        # Load the audio file
        audio_clip = AudioFileClip(audio_path)
        # Set the audio of the video clip to the loaded audio
        final_clip = video_clip.set_audio(audio_clip)
        # Export the final video with the new audio
        final_clip.write_videofile(output_path, codec='libx264', audio_codec='aac')
        print("Successfully created a video at "+ output_path)

    def embeded_cc(self,file_path):
        video_folder = os.path.join(file_path, "temp")
        output_folder=os.path.join(file_path,"output")
        video_path = os.path.join(video_folder,self.procID+".mp4")
        subtitle_path=os.path.join(output_folder, "sub-presetation.en.srt")
        output_path=os.path.join(output_folder,self.procID+".mp4")
        try:
        # Construct the ffmpeg command to embed the soft subtitle
            command = [
                'ffmpeg',
                '-i', video_path,  # Input video
                '-i', subtitle_path,  # Input subtitle
                '-c', 'copy',  # Copy video and audio without re-encoding
                '-c:s', 'mov_text',  # Codec for subtitles (mov_text for MP4 container)
                output_path
            ]
            # Execute the command using subprocess
            subprocess.run(command, check=True)
            print(f"Subtitle embedded successfully! Output saved to: {output_path}")
        except subprocess.CalledProcessError as e:
            print(f"Error: {e}")
            print("Failed to embed subtitle into video.")

    def clearUp(self,file_path):
        #delete all temporaty folders
        try:
            shutil.rmtree(os.path.join(file_path, "To_Mp4"))
            shutil.rmtree(os.path.join(file_path, "To_Wavs"))
            shutil.rmtree(os.path.join(file_path, "To_Pngs"))
            shutil.rmtree(os.path.join(file_path, "temp"))
            print("Successfully delete all redundant files")
        except:
            pass

def main(ID, voice_model):
    # print("Current working directory:", os.getcwd()) 
    # print("Python version:", sys.version)
    # print("Python executable:", sys.executable)
    current_dir=os.path.join(os.getcwd(),'media')
    cur_request=presentoai(ID,voice_model)
    procID=cur_request.procID
    voice_model=cur_request.voice_model
    print(procID)
    print(voice_model)
    folder_path=os.path.join(current_dir,procID)
    # print(folder_path)
    output_path=[]
    for folder in os.listdir(folder_path):
        folder_path=os.path.join(folder_path,folder)
        cur_request.clearUp(folder_path)
        cur_request.update_process_percentage(10)
        #create images from pptx
        cur_request.getImg(folder_path)
        cur_request.update_process_percentage(15) 
        #create mp3 files from speaker note
        cur_request.getWav(folder_path, voice_model)
        cur_request.update_process_percentage(35)
        #generate clip (mp4)
        cur_request.getClip(folder_path)
        cur_request.update_process_percentage(50)
        #combine all mp3 files 
        cur_request.combineAudio(folder_path)
        cur_request.getCC(folder_path)
        cur_request.update_process_percentage(75)
        #combine the mp3 and mp4 file
        cur_request.video_wsound(folder_path)
        cur_request.update_process_percentage(90)
        cur_request.embeded_cc(folder_path)
        cur_request.update_process_percentage(95)
        #make sure remove all unnecessary folder
        cur_request.clearUp(folder_path)
        cur_request.update_process_percentage(100)
        output_path.append(os.path.join(folder_path,"output",folder+".mp4")) 
        return(output_path)

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Process some IDs and voice model.")
    parser.add_argument('process_id', type=str, help='The process ID to be used')
    parser.add_argument('voice_model', type=str, help='The voice model to be used')
    args = parser.parse_args()
    main(args.process_id, args.voice_model)

#main(process_id="02bce7ee95924b29b6098dff016b4748",voice_model=1)