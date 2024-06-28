from pptx import Presentation

mypath=r"C:\Users\Thanh Lu\OneDrive\tts\sample.pptx"
def getText (my_path):
    ppt=Presentation(my_path)
    notes = []
    for slide in ppt.slides:
        textNote = slide.notes_slide.notes_text_frame.text
        notes.append(textNote) 
    return notes

